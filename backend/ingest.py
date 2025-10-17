"""
Document ingestion module for PDF, DOCX, PPTX, images, and text files.
Handles extraction, chunking, and metadata tracking.
"""

import os
import io
from typing import List, Dict, Any, Tuple
from dataclasses import dataclass
import fitz  # PyMuPDF
import pdfplumber
from docx import Document
from pptx import Presentation
from PIL import Image
import tiktoken


@dataclass
class DocumentChunk:
    """Represents a chunk of document content with metadata."""
    chunk_id: str
    filename: str
    page: int
    chunk_text: str
    token_count: int
    chunk_index: int
    file_type: str
    priority: str = "normal"  # rubric, slides, textbook, normal
    trusted: bool = True
    metadata: Dict[str, Any] = None


class DocumentIngestor:
    """Handles ingestion of various document formats."""
    
    def __init__(self, chunk_size: int = 600, chunk_overlap: int = 100):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.tokenizer = tiktoken.get_encoding("cl100k_base")
    
    def ingest_file(self, filepath: str, priority: str = "normal", trusted: bool = True) -> List[DocumentChunk]:
        """
        Ingest a file and return chunks with metadata.
        
        Args:
            filepath: Path to the file
            priority: Priority level (rubric, slides, textbook, normal)
            trusted: Whether the source is trusted
            
        Returns:
            List of DocumentChunk objects
        """
        filename = os.path.basename(filepath)
        ext = os.path.splitext(filename)[1].lower()
        
        if ext == '.pdf':
            return self._ingest_pdf(filepath, filename, priority, trusted)
        elif ext == '.docx':
            return self._ingest_docx(filepath, filename, priority, trusted)
        elif ext == '.pptx':
            return self._ingest_pptx(filepath, filename, priority, trusted)
        elif ext in ['.txt', '.text']:
            return self._ingest_text(filepath, filename, priority, trusted)
        elif ext in ['.png', '.jpg', '.jpeg']:
            return self._ingest_image(filepath, filename, priority, trusted)
        else:
            raise ValueError(f"Unsupported file type: {ext}")
    
    def _ingest_pdf(self, filepath: str, filename: str, priority: str, trusted: bool) -> List[DocumentChunk]:
        """Extract text from PDF using PyMuPDF."""
        chunks = []
        doc = fitz.open(filepath)
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text("text")
            
            # Also extract text blocks with positions for later highlighting
            blocks = page.get_text("dict")["blocks"]
            
            if text.strip():
                page_chunks = self._chunk_text(
                    text, 
                    filename, 
                    page_num + 1, 
                    "pdf", 
                    priority, 
                    trusted,
                    {"blocks": blocks}
                )
                chunks.extend(page_chunks)
        
        doc.close()
        return chunks
    
    def _ingest_docx(self, filepath: str, filename: str, priority: str, trusted: bool) -> List[DocumentChunk]:
        """Extract text from DOCX."""
        doc = Document(filepath)
        full_text = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
        
        text = "\n".join(full_text)
        return self._chunk_text(text, filename, 1, "docx", priority, trusted)
    
    def _ingest_pptx(self, filepath: str, filename: str, priority: str, trusted: bool) -> List[DocumentChunk]:
        """Extract text from PPTX."""
        prs = Presentation(filepath)
        chunks = []
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                text = "\n".join(slide_text)
                slide_chunks = self._chunk_text(
                    text, 
                    filename, 
                    slide_num, 
                    "pptx", 
                    priority, 
                    trusted
                )
                chunks.extend(slide_chunks)
        
        return chunks
    
    def _ingest_text(self, filepath: str, filename: str, priority: str, trusted: bool) -> List[DocumentChunk]:
        """Ingest plain text file."""
        with open(filepath, 'r', encoding='utf-8') as f:
            text = f.read()
        
        return self._chunk_text(text, filename, 1, "txt", priority, trusted)
    
    def _ingest_image(self, filepath: str, filename: str, priority: str, trusted: bool) -> List[DocumentChunk]:
        """Ingest image - placeholder for OCR."""
        # This will be processed by ocr_math.py
        return [DocumentChunk(
            chunk_id=f"{filename}_img_1",
            filename=filename,
            page=1,
            chunk_text=f"[IMAGE: {filename} - requires OCR processing]",
            token_count=0,
            chunk_index=0,
            file_type="image",
            priority=priority,
            trusted=trusted,
            metadata={"filepath": filepath}
        )]
    
    def _chunk_text(
        self, 
        text: str, 
        filename: str, 
        page: int, 
        file_type: str, 
        priority: str, 
        trusted: bool,
        metadata: Dict = None
    ) -> List[DocumentChunk]:
        """
        Split text into chunks based on token count.
        
        Args:
            text: Text to chunk
            filename: Source filename
            page: Page number
            file_type: Type of file
            priority: Priority level
            trusted: Whether source is trusted
            metadata: Additional metadata
            
        Returns:
            List of DocumentChunk objects
        """
        # Tokenize the text
        tokens = self.tokenizer.encode(text)
        chunks = []
        
        # Create overlapping chunks
        start = 0
        chunk_idx = 0
        
        while start < len(tokens):
            end = min(start + self.chunk_size, len(tokens))
            chunk_tokens = tokens[start:end]
            chunk_text = self.tokenizer.decode(chunk_tokens)
            
            chunk_id = f"{filename}_p{page}_c{chunk_idx}"
            
            chunk = DocumentChunk(
                chunk_id=chunk_id,
                filename=filename,
                page=page,
                chunk_text=chunk_text.strip(),
                token_count=len(chunk_tokens),
                chunk_index=chunk_idx,
                file_type=file_type,
                priority=priority,
                trusted=trusted,
                metadata=metadata or {}
            )
            
            chunks.append(chunk)
            chunk_idx += 1
            
            # Move to next chunk with overlap
            start += (self.chunk_size - self.chunk_overlap)
        
        return chunks
    
    def get_page_image(self, pdf_path: str, page_num: int, dpi: int = 150) -> Image.Image:
        """
        Extract a page as an image for display/highlighting.
        
        Args:
            pdf_path: Path to PDF file
            page_num: Page number (1-indexed)
            dpi: Resolution for rendering
            
        Returns:
            PIL Image object
        """
        from pdf2image import convert_from_path
        
        images = convert_from_path(
            pdf_path,
            first_page=page_num,
            last_page=page_num,
            dpi=dpi
        )
        
        return images[0] if images else None
    
    def get_text_coordinates(self, pdf_path: str, page_num: int, search_text: str) -> List[Tuple[float, float, float, float]]:
        """
        Get bounding box coordinates for text in PDF.
        
        Args:
            pdf_path: Path to PDF
            page_num: Page number (1-indexed)
            search_text: Text to search for
            
        Returns:
            List of (x0, y0, x1, y1) bounding boxes
        """
        doc = fitz.open(pdf_path)
        page = doc[page_num - 1]
        
        # Search for text instances
        text_instances = page.search_for(search_text[:100])  # First 100 chars
        
        doc.close()
        return [(rect.x0, rect.y0, rect.x1, rect.y1) for rect in text_instances]


def extract_file_metadata(filepath: str) -> Dict[str, Any]:
    """
    Extract metadata from file.
    
    Args:
        filepath: Path to file
        
    Returns:
        Dictionary with metadata
    """
    filename = os.path.basename(filepath)
    ext = os.path.splitext(filename)[1].lower()
    file_size = os.path.getsize(filepath)
    
    metadata = {
        "filename": filename,
        "extension": ext,
        "size_bytes": file_size,
        "size_mb": round(file_size / (1024 * 1024), 2),
    }
    
    if ext == '.pdf':
        doc = fitz.open(filepath)
        metadata["page_count"] = len(doc)
        doc.close()
    elif ext == '.pptx':
        prs = Presentation(filepath)
        metadata["page_count"] = len(prs.slides)
    elif ext == '.docx':
        doc = Document(filepath)
        metadata["page_count"] = len(doc.paragraphs)
    else:
        metadata["page_count"] = 1
    
    return metadata
